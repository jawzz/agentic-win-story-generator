"""
Win Story Generator - PPTX Engine v2
Builds a single-slide agentic use case deck using python-pptx.

v2 changes:
- Role colors: BOT=orange, AGENT=teal, HUMAN=gold (swapped from v1)
- Default theme = light
- Removed top-left orange dash (cleaner top zone)
- Tightened gap between Problem/Solution card headers and descriptions
- Proper vertical gap between Measured outcomes and Impact cards (no overlap)
- Attributable/Downstream accept either a string OR a list of {direction, text} dicts
- Fetches customer company logo via Clearbit when company domain can be inferred
"""
from pathlib import Path
import io
import urllib.request
import urllib.parse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Swapped: BOT=orange, AGENT=teal, HUMAN=gold
DARK = {
    'BG':C('14222A'),'CARD_BG':C('1E3038'),'CARD_BG_2':C('253A45'),'CARD_BG_3':C('2E4552'),
    'DIVIDER':C('2A3F4A'),'ORANGE':C('FA4616'),'ORANGE_CARD':C('7A2E18'),
    'TEAL':C('0BA2B3'),'TEAL_CARD':C('0E5C68'),'GOLD':C('DA9100'),
    'GREEN':C('5BBE82'),'RED':C('E53E3E'),'DEEP_BLUE':C('1E6482'),
    'WHITE':C('FFFFFF'),'CREAM':C('FFE8DC'),'TEXT':C('FFFFFF'),
    'TEXT_MUTED':C('8AABB5'),'TEXT_DIM':C('5C7480'),
    'OUTCOME_FILL':C('FA4616'),
}
LIGHT = {
    'BG':C('F7F8FA'),'CARD_BG':C('FFFFFF'),'CARD_BG_2':C('EDEFF2'),'CARD_BG_3':C('FFFFFF'),
    'DIVIDER':C('D8DDE3'),'ORANGE':C('FA4616'),'ORANGE_CARD':C('0BA2B3'),
    'TEAL':C('0BA2B3'),'TEAL_CARD':C('1E6482'),'GOLD':C('DA9100'),
    'GREEN':C('5BBE82'),'RED':C('D43A2C'),'DEEP_BLUE':C('1E6482'),
    'WHITE':C('FFFFFF'),'CREAM':C('FFE8DC'),'TEXT':C('1A2330'),
    'TEXT_MUTED':C('5A6B78'),'TEXT_DIM':C('8A98A5'),
    'OUTCOME_FILL':C('FA4616'),
}

F = "Poppins"
BASE_DIR = Path(__file__).parent
LOGO_WHITE = BASE_DIR / 'static' / 'uipath_logo_2400.png'
LOGO_DARK = BASE_DIR / 'static' / 'uipath_logo_2400_dark.png'
BPMN_ROBOT = BASE_DIR / 'static' / 'icons' / 'bpmn_robot.png'
BPMN_AGENT = BASE_DIR / 'static' / 'icons' / 'bpmn_agent.png'
BPMN_PERSON = BASE_DIR / 'static' / 'icons' / 'bpmn_person.png'


# ---------- shape helpers ----------
def _rect(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _rect_bordered(s, x, y, w, h, fill, border, bw=0.5):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(bw)
    sh.shadow.inherit = False
    return sh


def _pill(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _chevron(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _text(s, x, y, w, h, text, *, size=12, bold=False, italic=False, color=None,
          align='left', anchor='top', tracking=None, font=F):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if anchor == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if anchor == 'bottom': tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    if align == 'right': p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    if color is not None: r.font.color.rgb = color
    if tracking is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking * 100)))
    return tb


def _mixed(s, x, y, w, h, parts, *, align='left', anchor='top'):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if anchor == 'middle': tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == 'center': p.alignment = PP_ALIGN.CENTER
    if align == 'right': p.alignment = PP_ALIGN.RIGHT
    for part in parts:
        r = p.add_run(); r.text = part['text']
        r.font.name = part.get('font', F)
        r.font.size = Pt(part.get('size', 12))
        r.font.bold = part.get('bold', False)
        r.font.italic = part.get('italic', False)
        if 'color' in part: r.font.color.rgb = part['color']
        if part.get('tracking') is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(int(part['tracking'] * 100)))
    return tb


# ---------- customer logo fetcher ----------
def _guess_domain(company):
    """Rough company name -> domain guess. Falls through common TLDs."""
    if not company: return None
    base = company.lower().strip()
    # Strip common suffixes
    for suf in [' inc', ' inc.', ' corporation', ' corp.', ' corp', ' llc', ' ltd',
                ' co.', ' company', ' plc', ' holdings', ' group', ',']:
        if base.endswith(suf):
            base = base[:-len(suf)].strip()
    base = base.replace(' & ', ' and ')
    domain_core = ''.join(c for c in base if c.isalnum() or c in '-')
    if not domain_core: return None
    return f'{domain_core}.com'


def _normalize_to_png(data):
    """Accept any image bytes (ICO, JPEG, PNG, etc.), return PNG bytes.
    python-pptx only embeds BMP/GIF/JPEG/PNG/TIFF/WMF — ICO fallbacks break it.
    We re-encode through PIL to PNG."""
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(data))
        # ICO files may contain multiple sizes — pick the largest
        if getattr(im, 'n_frames', 1) > 1 or im.format == 'ICO':
            try:
                sizes = im.info.get('sizes') or []
                if sizes:
                    largest = max(sizes, key=lambda s: s[0]*s[1])
                    im.size = largest
                    im.load()
            except Exception:
                pass
        # Convert mode to RGBA to preserve transparency
        if im.mode not in ('RGB', 'RGBA'):
            im = im.convert('RGBA')
        buf = io.BytesIO()
        im.save(buf, 'PNG')
        return buf.getvalue()
    except Exception:
        return None


def _fetch_customer_logo(company, timeout=4):
    """Fetch a customer logo image by guessing the company's domain.
    Tries Clearbit first (high-res SVG/PNG), then Google S2 favicon,
    then DuckDuckGo icon service as fallbacks. Matches behavior of the
    original Win Story Generator. Returns PNG image bytes or None."""
    domain = _guess_domain(company)
    if not domain:
        return None
    sources = [
        f'https://logo.clearbit.com/{domain}?size=512',
        f'https://logo.clearbit.com/{domain}?size=256',
        f'https://www.google.com/s2/favicons?domain={domain}&sz=256',
        f'https://www.google.com/s2/favicons?domain={domain}&sz=128',
        f'https://icons.duckduckgo.com/ip3/{domain}.ico',
    ]
    headers = {'User-Agent': 'Mozilla/5.0 WinStoryGenerator'}
    for url in sources:
        try:
            req = urllib.request.Request(url, headers=headers)
            with urllib.request.urlopen(req, timeout=timeout) as resp:
                if resp.status == 200:
                    data = resp.read()
                    if len(data) > 200:
                        # Always normalize to PNG so python-pptx accepts it
                        png = _normalize_to_png(data)
                        if png:
                            return png
        except Exception:
            continue
    return None



# ---------- UiPath + Customer partnership lockup ----------
def _draw_plus_circle(slide, cx, cy, diameter, fill_color, cross_color):
    """Draw a thick-outlined ring with a + cross inside (matches attached icon).
    Transparent inside of ring, black cross and black ring stroke.
    `fill_color` is unused (kept for signature compatibility). `cross_color` is
    used for BOTH the ring outline and the + bars."""
    r = diameter / 2
    # Outlined ring (no fill, thick black border)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx - r), Inches(cy - r),
                                  Inches(diameter), Inches(diameter))
    ring.fill.background()  # transparent inside
    ring.line.color.rgb = cross_color
    # Stroke thickness scales with diameter so the ring reads as thick
    stroke_pt = max(1.5, diameter * 5.5)  # ~0.34in -> ~2.7pt
    ring.line.width = Pt(stroke_pt)
    ring.shadow.inherit = False

    # + cross inside (two thick bars)
    bar_len = diameter * 0.44
    bar_thick = diameter * 0.17
    # horizontal bar
    hx = cx - bar_len / 2
    hy = cy - bar_thick / 2
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(hx), Inches(hy), Inches(bar_len), Inches(bar_thick))
    h.fill.solid(); h.fill.fore_color.rgb = cross_color
    h.line.fill.background(); h.shadow.inherit = False
    # vertical bar
    vx = cx - bar_thick / 2
    vy = cy - bar_len / 2
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(vx), Inches(vy), Inches(bar_thick), Inches(bar_len))
    v.fill.solid(); v.fill.fore_color.rgb = cross_color
    v.line.fill.background(); v.shadow.inherit = False


def _draw_tags(slide, T, data, divider_y):
    """Render small pill tags (INTERNAL, MAESTRO, EASY PROCESS) just above the divider on the right."""
    tags = []
    # Classification flag — only "internal" raises the INTERNAL USE pill.
    # "anonymize" handled separately: customer name swapped for industry-generic
    # label upstream, no pill needed.
    cls = (data.get('classification') or '').lower()
    if cls == 'internal' or data.get('internal'):
        tags.append(('FOR INTERNAL USE ONLY', T['RED']))
    # Maestro pill is rendered on the solution card now, not in the top-right row.
    # Easy process — explicit flag OR auto-detect (<=5 steps, no AGENT, no IXP)
    easy_flag = data.get('easy_process')
    if easy_flag is None:
        steps = data.get('steps') or []
        roles = []
        for s in steps:
            if isinstance(s, dict):
                roles.append(str(s.get('role','')).upper())
        easy_flag = (0 < len(steps) <= 5) and not any(r in ('AGENT', 'IXP') for r in roles)
    if easy_flag:
        tags.append(('EASY PROCESS', T['GREEN']))

    if not tags:
        return
    # Render as a horizontal row, right-aligned, just above the divider.
    # Pills are sized per-tag — INTERNAL pill is bigger to be unmissable.
    gap = 0.12
    # Bumped up from divider_y - 0.10 so pill row sits a touch higher on the page
    y_anchor_bottom = divider_y - 0.20
    char_w_big = 0.080
    char_w_small = 0.075
    pill_pad_big = 0.22
    pill_pad_small = 0.18
    sized = []
    for (label, color) in tags:
        if 'INTERNAL' in label:
            h = 0.32
            w = pill_pad_big*2 + len(label) * char_w_big
            font_size = 10
        else:
            h = 0.28
            w = pill_pad_small*2 + len(label) * char_w_small
            font_size = 9
        sized.append((label, color, w, h, font_size))
    total = sum(s[2] for s in sized) + gap * (len(sized) - 1)
    right_edge = 13.333 - 0.5
    x = right_edge - total
    for (label, color, w, h, font_size) in sized:
        # Bottom-align all pills so they share the same baseline
        y = y_anchor_bottom - h
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
        s.adjustments[0] = 0.5
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        _text(slide, x, y, w, h, label, size=font_size, bold=True,
              color=T['WHITE'], align='center', anchor='middle', tracking=1.2)
        x += w + gap


def _draw_partnership_lockup(slide, uipath_logo_path, customer_logo_bytes, company_name,
                              *, text_color, muted, plus_color):
    """UiPath logo is ALWAYS fixed in the top-right corner (same spot every time).
    Customer element (logo or name) sits to the LEFT of UiPath, with a circled +
    between them. If no company_name, only the UiPath logo renders."""
    import tempfile, os as _os
    from PIL import Image as _PIL

    # ---- UiPath logo: FIXED position, top-right ----
    try:
        uim = _PIL.open(uipath_logo_path); uiw, uih = uim.size
    except Exception:
        uiw, uih = 3, 1
    ui_h = 0.40
    ui_w = ui_h * (uiw / uih if uih else 3.0)
    ui_right = 12.90
    ui_top = 0.30
    ui_left = ui_right - ui_w
    ui_center_y = ui_top + ui_h / 2
    slide.shapes.add_picture(uipath_logo_path, Inches(ui_left), Inches(ui_top),
                             width=Inches(ui_w), height=Inches(ui_h))

    if not company_name:
        return

    # ---- Plus-in-circle sits to the LEFT of the UiPath logo ----
    plus_diameter = 0.22
    plus_gap = 0.22  # symmetric gap on both sides of the + icon
    plus_cx = ui_left - plus_gap - plus_diameter / 2
    plus_cy = ui_center_y + 0.05  # small nudge below geometric center to hit UiPath wordmark
    # Icon = black (or theme text color) — transparent inside of ring
    icon_color = text_color
    _draw_plus_circle(slide, plus_cx, plus_cy, plus_diameter, None, icon_color)

    # ---- Customer element to the LEFT of the plus ----
    cust_right = plus_cx - plus_diameter / 2 - plus_gap
    if customer_logo_bytes:
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        tmp.write(customer_logo_bytes); tmp.close()
        try:
            try:
                cim = _PIL.open(tmp.name); cw, ch = cim.size
            except Exception:
                cw, ch = 1, 1
            # Preserve aspect ratio, fit into max area
            c_max_h = 0.50
            c_max_w = 2.20
            ar_c = cw / ch if ch else 1.0
            c_h = c_max_h
            c_w = c_h * ar_c
            if c_w > c_max_w:
                c_w = c_max_w
                c_h = c_w / ar_c
            cust_x = cust_right - c_w
            cust_y = ui_center_y - c_h / 2 + 0.05
            slide.shapes.add_picture(tmp.name, Inches(cust_x), Inches(cust_y),
                                     width=Inches(c_w), height=Inches(c_h))
        finally:
            try: _os.unlink(tmp.name)
            except Exception: pass
    else:
        # Customer as bold text, right-aligned to cust_right
        name_w = min(2.4, max(0.8, 0.14 * len(company_name) + 0.3))
        name_x = cust_right - name_w
        name_y = ui_center_y - 0.10
        _text(slide, name_x, name_y, name_w, 0.38, company_name,
              size=15, bold=True, color=text_color, align='right', anchor='middle')


# ---------- main slide builder ----------
def _build_slide(slide, *, theme, data):
    T = theme
    is_light = theme is LIGHT
    logo_path = str(LOGO_DARK if is_light else LOGO_WHITE)

    title_str = data.get('title', '')
    is_short_title = len(title_str) < 55  # slightly relaxed threshold
    if is_short_title:
        divider_y = 1.52
        row_y = 1.65
        row_h = 1.95
        orch_ch = 1.40
        step_h = 0.90
    else:
        divider_y = 1.92
        row_y = 2.05
        row_h = 1.85
        orch_ch = 1.30
        step_h = 0.80
    gap = 0.30
    col_w = (12.33 - gap) / 2

    # Background ONLY — no orange dash
    _rect(slide, 0, 0, 13.333, 7.5, T['BG'])

    # Breadcrumb + UiPath logo (logo height bumped a touch for readability)
    bc = data.get('breadcrumb', ('Industry', 'Function', 'Use case name'))
    if isinstance(bc, (list, tuple)) and len(bc) >= 3:
        bc0, bc1, bc2 = bc[0], bc[1], bc[2]
    else:
        bc0, bc1, bc2 = 'Industry', 'Function', 'Use case name'
    _mixed(slide, 0.5, 0.30, 9.0, 0.28, [
        {'text': bc0, 'size': 10, 'color': T['TEXT_MUTED'], 'bold': True},
        {'text': '   /   ', 'size': 10, 'color': T['TEXT_DIM']},
        {'text': bc1, 'size': 10, 'color': T['TEXT_MUTED'], 'bold': True},
        {'text': '   /   ', 'size': 10, 'color': T['TEXT_DIM']},
        {'text': bc2, 'size': 10, 'color': T['TEXT'], 'bold': True},
    ])
    # UiPath logo + customer logo lockup placed below (uses _draw_logo_lockup)

    # Title + subtitle
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.60), Inches(9.0), Inches(1.20))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.space_after = Pt(4)
    r1 = p1.add_run(); r1.text = title_str
    r1.font.name = F; r1.font.size = Pt(26); r1.font.bold = True
    r1.font.color.rgb = T['TEXT']
    subtitle = data.get('subtitle', '')
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = F; r2.font.size = Pt(13)
        r2.font.color.rgb = T['TEXT_MUTED']

    # --- UiPath + Customer lockup in top-right ---
    anonymize = bool(data.get('anonymize')) or (data.get('classification') or '').lower() == 'anonymize'
    company_name = data.get('company', '')
    if anonymize:
        # Replace customer name with industry-generic label, no logo
        ind = bc0 if bc0 and bc0 != 'Industry' else 'Anonymous'
        display_company = f"{ind} Customer"
        logo_bytes = None
    else:
        display_company = company_name
        logo_bytes = _fetch_customer_logo(company_name) if company_name else None
    _draw_partnership_lockup(slide, logo_path, logo_bytes, display_company,
                             text_color=T['TEXT'], muted=T['TEXT_MUTED'],
                             plus_color=T['TEXT_MUTED'])

    # --- Tag pills (Internal / Maestro / Easy Process) ---
    _draw_tags(slide, T, data, divider_y)

    _rect(slide, 0.5, divider_y, 12.33, 0.015, T['DIVIDER'])

    # PROBLEM card
    px = 0.5
    _rect(slide, px, row_y, col_w, row_h, T['ORANGE_CARD'])
    # Header
    _text(slide, px+0.35, row_y+0.18, col_w-0.7, 0.38,
          'The problem', size=18, bold=True, color=T['WHITE'])

    problem_desc = (data.get('problem_desc') or '').strip()
    problem_stats = data.get('problem_stats') or []
    norm_stats = []
    for st in problem_stats:
        if isinstance(st, dict):
            val = (st.get('value') or st.get('number') or '').strip()
            lbl = (st.get('label') or '').strip()
            if val or lbl:
                norm_stats.append((val, lbl))
        elif isinstance(st, (list, tuple)) and len(st) >= 2:
            norm_stats.append((str(st[0]), str(st[1])))

    # Description tightened up further (was 0.62, now 0.50 — closer to header)
    desc_top_prob = row_y + 0.50
    if norm_stats:
        # Description gets less room; stats below
        desc_h_prob = 0.48 if len(problem_desc) < 180 else 0.72
        # Auto-shrink desc size if very long
        desc_size_prob = 10
        if len(problem_desc) > 220: desc_size_prob = 9
        _text(slide, px+0.35, desc_top_prob, col_w-0.7, desc_h_prob,
              problem_desc, size=desc_size_prob, color=T['WHITE'])
        stats_area_h = 0.58
        stats_y = row_y + row_h - 0.14 - stats_area_h
        sa_w = col_w - 0.7
        sc = min(len(norm_stats), 4)
        sw = (sa_w - 0.15*(sc-1)) / sc
        for i, (n, l) in enumerate(norm_stats[:sc]):
            sx = px + 0.35 + i*(sw+0.15)
            # Drop to 18 when value > 6 chars (e.g. "$1.5M+", "21 days") so it fits
            stat_size = 18 if len(n) > 6 else 22
            _text(slide, sx, stats_y, sw, 0.34, n, size=stat_size, bold=True, color=T['WHITE'])
            _text(slide, sx, stats_y+0.36, sw, 0.26, l, size=9, color=T['CREAM'])
    else:
        # No stats — use full lower space for description
        desc_size_prob = 13 if len(problem_desc) < 200 else (12 if len(problem_desc) < 320 else 11)
        _text(slide, px+0.35, desc_top_prob, col_w-0.7, row_h-0.78,
              problem_desc, size=desc_size_prob, color=T['WHITE'])

    # SOLUTION card
    sx0 = px + col_w + gap
    _rect(slide, sx0, row_y, col_w, row_h, T['TEAL_CARD'])
    _text(slide, sx0+0.35, row_y+0.18, col_w-0.7, 0.38,
          'The solution', size=18, bold=True, color=T['WHITE'])

    pills = data.get('capabilities') or []
    max_pill_area_w = col_w - 0.7
    rows = [[]]; row_w = 0.0
    pill_h = 0.28; pill_row_gap = 0.06
    for L in pills:
        w = 0.22 + len(L)*0.075
        if rows[-1] and row_w + w > max_pill_area_w:
            rows.append([]); row_w = 0.0
        rows[-1].append((L, w)); row_w += w + 0.08
    n_rows = len(rows) if pills else 0
    pill_block_h = n_rows*pill_h + max(0, n_rows-1)*pill_row_gap if pills else 0
    bottom_margin = 0.10 if n_rows >= 2 else 0.20
    pill_block_bottom = row_y + row_h - bottom_margin
    pill_block_top = pill_block_bottom - pill_block_h

    desc_top_sol = row_y + 0.50  # tightened from 0.62
    solution_desc = (data.get('solution_desc') or '').strip()
    if pills:
        desc_h_sol = max(0.30, pill_block_top - 0.12 - desc_top_sol)
        # More room when fewer pills (single row)
        desc_size_sol = 11 if n_rows == 1 else 10
        if len(solution_desc) > 260: desc_size_sol = min(desc_size_sol, 9)
    else:
        desc_h_sol = row_h - 0.78
        desc_size_sol = 13 if len(solution_desc) < 200 else (12 if len(solution_desc) < 320 else 11)

    _text(slide, sx0+0.35, desc_top_sol, col_w-0.7, desc_h_sol,
          solution_desc, size=desc_size_sol, color=T['WHITE'])

    if pills:
        pill_fill = T['WHITE']
        pill_text = C('0E5C68')
        for ri, row in enumerate(rows):
            py_row = pill_block_top + ri*(pill_h + pill_row_gap)
            pxp = sx0 + 0.35
            for L, w in row:
                _pill(slide, pxp, py_row, w, pill_h, pill_fill)
                _text(slide, pxp, py_row, w, pill_h, L,
                      size=9, bold=True, color=pill_text, align='center', anchor='middle')
                pxp += w + 0.08

    # What the automation does
    cy = row_y + row_h + 0.10
    ch = orch_ch
    _rect(slide, 0.5, cy, 12.33, ch, T['CARD_BG_2'])
    _text(slide, 0.7, cy+0.12, 8.0, 0.30,
          'What the automation does', size=15, bold=True, color=T['TEXT'])

    # MAESTRO pill in the top-right of "What the automation does" card, if Maestro is in play
    _caps_for_maestro = [str(c).lower() for c in (data.get('capabilities') or [])]
    if data.get('maestro') or any('maestro' in c for c in _caps_for_maestro):
        m_label = 'MAESTRO'
        m_w = 0.18*2 + len(m_label)*0.075
        m_h = 0.28
        m_x = 0.5 + 12.33 - m_w - 0.30
        # Nudged up slightly (was cy + 0.13)
        m_y = cy + 0.09
        m_pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(m_x), Inches(m_y), Inches(m_w), Inches(m_h))
        m_pill.adjustments[0] = 0.5
        m_pill.fill.solid(); m_pill.fill.fore_color.rgb = T['DEEP_BLUE']
        m_pill.line.fill.background()
        m_pill.shadow.inherit = False
        _text(slide, m_x, m_y, m_w, m_h, m_label, size=9, bold=True,
              color=T['WHITE'], align='center', anchor='middle', tracking=1.2)

    steps = data.get('steps') or []
    norm_steps = []
    for st in steps:
        if isinstance(st, dict):
            role = (st.get('role') or st.get('type') or 'AGENT').upper()
            if role in ('ROBOT',): role = 'BOT'
            if role in ('PERSON',): role = 'HUMAN'
            # Accept AGENT, BOT, HUMAN, IXP — anything else falls back to AGENT
            if role not in ('AGENT', 'BOT', 'HUMAN', 'IXP'):
                role = 'AGENT'
            desc = (st.get('description') or st.get('name') or '').strip()
            norm_steps.append((role, desc))
        elif isinstance(st, (list, tuple)) and len(st) >= 2:
            r = str(st[0]).upper()
            if r not in ('AGENT', 'BOT', 'HUMAN', 'IXP'):
                r = 'AGENT'
            norm_steps.append((r, str(st[1])))

    # SWAPPED COLORS: BOT=orange, AGENT=teal, HUMAN=gold, IXP=green
    role_c = {'BOT': T['ORANGE'], 'AGENT': T['TEAL'], 'HUMAN': T['GOLD'], 'IXP': T['GREEN']}
    n_steps = len(norm_steps)

    if n_steps > 0:
        if n_steps <= 5:
            arrow_w, gap_s, show_num = 0.22, 0.10, True
        elif n_steps <= 6:
            arrow_w, gap_s, show_num = 0.18, 0.08, True
        elif n_steps <= 7:
            arrow_w, gap_s, show_num = 0.16, 0.07, True
        else:
            arrow_w, gap_s, show_num = 0.14, 0.06, False

        step_y = cy + 0.45
        inner_pad = 0.30
        avail_w = 12.33 - 2*inner_pad
        tile_w = (avail_w - arrow_w*(n_steps-1) - gap_s*(n_steps-1)) / n_steps

        def desc_size_for_step(desc, tw):
            avail_text_w = tw - 0.28
            chars_at_10 = max(1, int(avail_text_w / 0.08))
            chars_at_9 = max(1, int(avail_text_w / 0.072))
            if len(desc) <= chars_at_10 * 3: return 10
            if len(desc) <= chars_at_9 * 3: return 9
            return 8

        x = 0.5 + inner_pad
        for i, (r, desc) in enumerate(norm_steps):
            if is_light:
                _rect_bordered(slide, x, step_y, tile_w, step_h, T['CARD_BG_3'], T['DIVIDER'], 0.75)
            else:
                _rect(slide, x, step_y, tile_w, step_h, T['CARD_BG_3'])
            _rect(slide, x, step_y, tile_w, 0.05, role_c[r])
            label = f"{i+1:02d}  {r}" if show_num else r
            _text(slide, x+0.14, step_y+0.12, tile_w-0.24, 0.22,
                  label, size=9, bold=True, color=role_c[r], tracking=1.5)
            # Role icon in top-right of tile (icons dropped ~5px from prior position)
            icon_size = 0.22
            icon_x = x + tile_w - icon_size - 0.02
            icon_y = step_y + 0.07  # was 0.02 — pushed down 5px to better center
            if r == 'IXP':
                # Green circle background with an outlined (no-fill) document inside,
                # ~80% of prior size so the green shows through clearly.
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                                Inches(icon_x), Inches(icon_y),
                                                Inches(icon_size), Inches(icon_size))
                circle.fill.solid(); circle.fill.fore_color.rgb = role_c[r]
                circle.line.fill.background()
                circle.shadow.inherit = False
                # Document shape: outline only (no fill), 80% of previous size, centered
                doc_w = icon_size * 0.55 * 0.80
                doc_h = icon_size * 0.62 * 0.80
                doc_x = icon_x + (icon_size - doc_w) / 2
                doc_y = icon_y + (icon_size - doc_h) / 2
                doc = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DOCUMENT,
                                              Inches(doc_x), Inches(doc_y),
                                              Inches(doc_w), Inches(doc_h))
                doc.fill.background()  # transparent — green shows through
                doc.line.color.rgb = T['WHITE']
                doc.line.width = Pt(0.25)
                doc.shadow.inherit = False
            else:
                icon_map = {'BOT': BPMN_ROBOT, 'AGENT': BPMN_AGENT, 'HUMAN': BPMN_PERSON}
                icon_path = icon_map.get(r)
                if icon_path and icon_path.exists():
                    slide.shapes.add_picture(str(icon_path), Inches(icon_x), Inches(icon_y),
                                              width=Inches(icon_size), height=Inches(icon_size))
            ds = desc_size_for_step(desc, tile_w)
            _text(slide, x+0.14, step_y+0.34, tile_w-0.24, step_h-0.38,
                  desc, size=ds, bold=True, color=T['TEXT'])
            x += tile_w
            if i < n_steps - 1:
                ax = x + gap_s/3
                ay = step_y + step_h/2 - arrow_w/2
                _chevron(slide, ax, ay, arrow_w, arrow_w, T['TEXT_MUTED'])
                x += arrow_w + gap_s*2/3

    # Measured outcomes (no caption, orange tiles)
    oy = cy + ch + 0.10
    _text(slide, 0.5, oy, 8.0, 0.30,
          'Measured outcomes', size=15, bold=True, color=T['TEXT'])

    outcomes = data.get('outcomes') or []
    norm_outcomes = []
    for o in outcomes:
        if isinstance(o, dict):
            val = (o.get('value') or o.get('number') or '').strip()
            lbl = (o.get('label') or '').strip()
            if val or lbl:
                norm_outcomes.append((val, lbl))
        elif isinstance(o, (list, tuple)) and len(o) >= 2:
            norm_outcomes.append((str(o[0]), str(o[1])))

    if norm_outcomes:
        nt = len(norm_outcomes)
        tile_y = oy + 0.30   # tightened from 0.40 — closer to header
        tile_h_out = 0.94    # taller (was 0.78) so 2-line labels fit
        tg = 0.20
        tw2 = (12.33 - tg*(nt-1)) / nt
        for i, (num, label) in enumerate(norm_outcomes):
            tx = 0.5 + i*(tw2 + tg)
            _rect(slide, tx, tile_y, tw2, tile_h_out, T['OUTCOME_FILL'])
            _text(slide, tx+0.32, tile_y+0.08, tw2-0.5, 0.42,
                  num, size=30, bold=True, color=T['WHITE'], anchor='middle')
            _text(slide, tx+0.32, tile_y+0.52, tw2-0.5, 0.40,
                  label, size=11, color=T['CREAM'], anchor='middle')
        outcomes_bottom = tile_y + tile_h_out
    else:
        outcomes_bottom = oy + 0.30

    def _format_impact(val):
        if isinstance(val, list):
            parts = []
            for item in val:
                if isinstance(item, dict):
                    d = (item.get('direction') or item.get('dir') or '').lower()
                    t = (item.get('text') or item.get('label') or '').strip()
                    if not t: continue
                    if d in ('up', 'increase', 'improved', 'better'): arrow = '\u2191'
                    elif d in ('down', 'decrease', 'reduced', 'worse'): arrow = '\u2193'
                    else: arrow = ''
                    parts.append(f'{arrow} {t}'.strip())
            return '   \u00b7   '.join(parts)
        return (val or '').strip() if val else ''

    attributable = _format_impact(data.get('attributable'))
    downstream = _format_impact(data.get('downstream'))

    ay = outcomes_bottom + 0.10  # tightened from 0.18
    ah = 0.78
    acol_w = (12.33 - gap) / 2

    if attributable or downstream:
        if attributable:
            if is_light:
                _rect_bordered(slide, 0.5, ay, acol_w, ah, T['CARD_BG'], T['DIVIDER'], 0.75)
            else:
                _rect(slide, 0.5, ay, acol_w, ah, T['CARD_BG'])
            _mixed(slide, 0.5+0.2, ay+0.10, acol_w-0.4, 0.26, [
                {'text': 'Attributable impact', 'size': 12, 'bold': True, 'color': T['TEAL']},
                {'text': '   directly moved, not yet quantified',
                 'size': 9, 'italic': True, 'color': T['TEXT_DIM']},
            ])
            metric_size = 9 if len(attributable) > 80 else 10
            _text(slide, 0.5+0.2, ay+0.36, acol_w-0.4, 0.38,
                  attributable, size=metric_size, color=T['TEXT'])

        if downstream:
            dx = 0.5 + acol_w + gap
            if is_light:
                _rect_bordered(slide, dx, ay, acol_w, ah, T['CARD_BG'], T['DIVIDER'], 0.75)
            else:
                _rect(slide, dx, ay, acol_w, ah, T['CARD_BG'])
            _mixed(slide, dx+0.2, ay+0.10, acol_w-0.4, 0.26, [
                {'text': 'Downstream impact', 'size': 12, 'bold': True, 'color': T['GOLD']},
                {'text': '   second-order effects',
                 'size': 9, 'italic': True, 'color': T['TEXT_DIM']},
            ])
            metric_size = 9 if len(downstream) > 80 else 10
            _text(slide, dx+0.2, ay+0.36, acol_w-0.4, 0.38,
                  downstream, size=metric_size, color=T['TEXT'])


def build_pptx(data, template_path=None):
    """Builds a single-slide agentic use case PPTX.

    Data format:
        - breadcrumb: [industry, function, useCase]
        - title, subtitle, company: strings
        - problem_desc: string
        - problem_stats: [{value, label}, ...]  (optional, 0-4)
        - solution_desc: string
        - capabilities: [string, ...]
        - steps: [{role: AGENT|BOT|HUMAN, description}, ...]
        - outcomes: [{value, label}, ...]  (1-5)
        - attributable: string OR [{direction: up|down, text}, ...]  (optional)
        - downstream:   string OR [{direction: up|down, text}, ...]  (optional)
        - theme: "dark" or "light" (default "light")
    """
    theme_name = (data.get('theme') or 'light').lower()
    theme = DARK if theme_name == 'dark' else LIGHT

    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    slide = p.slides.add_slide(p.slide_layouts[6])
    _build_slide(slide, theme=theme, data=data)

    # Speaker notes — account team contact + classification metadata
    notes_lines = []
    if data.get('account_team'):
        notes_lines.append(f"Account team contact: {data['account_team']}")
    cls = (data.get('classification') or '').strip()
    if cls:
        notes_lines.append(f"Sharing classification: {cls}")
    if data.get('company'):
        notes_lines.append(f"Customer: {data['company']}")
    bc = data.get('breadcrumb') or []
    if len(bc) >= 2:
        notes_lines.append(f"Industry / Function: {bc[0]} / {bc[1]}")
    caps = data.get('capabilities') or []
    if caps:
        notes_lines.append(f"Capabilities: {', '.join(str(c) for c in caps)}")
    if notes_lines:
        try:
            slide.notes_slide.notes_text_frame.text = '\n'.join(notes_lines)
        except Exception:
            pass

    buf = io.BytesIO()
    p.save(buf)
    buf.seek(0)
    return buf.read(), 1
